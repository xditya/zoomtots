from pptx import Presentation

def extract_text_from_ppt(ppt_file_path):
    prs = Presentation(ppt_file_path)
    slides_text = []

    for slide in prs.slides:
        text = ''
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + '\n'
        slides_text.append(text.strip())

    return slides_text
